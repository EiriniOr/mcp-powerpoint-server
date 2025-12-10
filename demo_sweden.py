#!/usr/bin/env python3
"""
Demo: Life in Sweden Presentation
Showcases the PowerPoint MCP Server's 36 tools
"""

import asyncio
import sys
import os

# Add the current directory to the path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from server import presentations, Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

async def create_sweden_demo():
    """Create a comprehensive demo presentation about Life in Sweden"""

    filename = "life_in_sweden_demo.pptx"

    print("Creating 'Life in Sweden' demo presentation...")
    print("This showcases the PowerPoint MCP Server's capabilities\n")

    # 1. Create presentation with title slide
    print("1. Creating title slide...")
    prs = Presentation()
    presentations[filename] = prs

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use blank layout for custom design

    # Add Swedish flag-inspired background shapes
    from pptx.enum.shapes import MSO_SHAPE

    # Blue background (full slide)
    bg_rect = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg_rect.fill.solid()
    bg_rect.fill.fore_color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    bg_rect.line.fill.background()

    # Yellow accent stripe (horizontal)
    accent = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(3.25),
        Inches(10), Inches(1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(254, 204, 0)  # Swedish yellow
    accent.line.fill.background()

    # Decorative circles (Nordic design elements)
    circle_positions = [(0.5, 0.5), (9, 0.5), (0.5, 6.5), (9, 6.5)]
    for x, y in circle_positions:
        circle = title_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(254, 204, 0)
        circle.line.fill.background()

    # Main title
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "Life in Sweden"
    title_frame.paragraphs[0].font.size = Pt(72)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    from pptx.enum.text import PP_ALIGN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle with better styling
    subtitle_box = title_slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Nordic Adventure"
    subtitle_frame.paragraphs[0].font.size = Pt(36)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add demo credit
    p = subtitle_frame.add_paragraph()
    p.text = "\nPowered by PowerPoint MCP Server"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 2. Add agenda slide with visual boxes
    print("2. Adding agenda slide...")
    agenda_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light blue background
    agenda_slide.background.fill.solid()
    agenda_slide.background.fill.fore_color.rgb = RGBColor(240, 248, 255)

    # Title
    title_box = agenda_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📋 Agenda"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)

    # Agenda items with colorful boxes
    agenda_items = [
        ("Swedish Culture & Traditions", RGBColor(0, 106, 167)),
        ("Cost of Living Comparison", RGBColor(30, 130, 180)),
        ("Climate Throughout the Year", RGBColor(60, 150, 195)),
        ("Popular Cities & Destinations", RGBColor(90, 170, 210)),
        ("Work-Life Balance & Process", RGBColor(120, 190, 225))
    ]

    y_start = 1.7
    for i, (item, color) in enumerate(agenda_items):
        # Number circle
        num_circle = agenda_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.8), Inches(y_start + i * 0.95),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = RGBColor(254, 204, 0)
        num_circle.line.fill.background()
        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        num_text.paragraphs[0].font.size = Pt(20)
        num_text.paragraphs[0].font.bold = True
        num_text.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
        num_text.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Item box with gradient-style colors
        item_box = agenda_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_start + i * 0.95),
            Inches(7.5), Inches(0.7)
        )
        item_box.fill.solid()
        item_box.fill.fore_color.rgb = color
        item_box.line.fill.background()
        item_text = item_box.text_frame
        item_text.text = item
        item_text.paragraphs[0].font.size = Pt(18)
        item_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        item_text.vertical_anchor = 1  # Middle

    # 3. Add section break with more visual interest
    print("3. Adding section break...")
    blank_slide_layout = prs.slide_layouts[6]
    section_slide = prs.slides.add_slide(blank_slide_layout)
    section_slide.background.fill.solid()
    section_slide.background.fill.fore_color.rgb = RGBColor(0, 106, 167)  # Swedish blue

    # Add decorative hexagons (Nordic pattern)
    hex_positions = [
        (0.5, 1), (2, 0.5), (3.5, 1), (5, 0.5), (6.5, 1), (8, 0.5),
        (0.5, 5.5), (2, 6), (3.5, 5.5), (5, 6), (6.5, 5.5), (8, 6)
    ]
    for x, y in hex_positions:
        hex_shape = section_slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            Inches(x), Inches(y),
            Inches(1), Inches(1)
        )
        hex_shape.fill.solid()
        hex_shape.fill.fore_color.rgb = RGBColor(254, 204, 0)
        hex_shape.fill.transparency = 0.7
        hex_shape.line.fill.background()

    # Large title in center
    title_box = section_slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.8))
    title_frame = title_box.text_frame
    title_frame.text = "Swedish Culture"
    title_frame.paragraphs[0].font.size = Pt(66)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = section_slide.shapes.add_textbox(Inches(2), Inches(4.8), Inches(6), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Traditions • Values • Lifestyle"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(254, 204, 0)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 4. Add comparison slide: Sweden vs Other Countries (enhanced)
    print("4. Adding comparison slide...")
    comparison_slide = prs.slides.add_slide(prs.slide_layouts[6])
    comparison_slide.background.fill.solid()
    comparison_slide.background.fill.fore_color.rgb = RGBColor(245, 250, 255)

    # Title with icon
    title_box = comparison_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "⚖️ Sweden vs European Neighbors"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left side - Sweden with rounded rectangle
    left_box = comparison_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.5), Inches(4), Inches(5.2)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(0, 106, 167)
    left_box.line.width = Pt(3)
    left_box.line.color.rgb = RGBColor(254, 204, 0)

    left_frame = left_box.text_frame
    left_frame.text = "Sweden 🇸🇪"
    left_frame.paragraphs[0].font.size = Pt(28)
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.color.rgb = RGBColor(254, 204, 0)
    left_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    left_items = [
        "👥 Population: ~10.5M",
        "⏰ Work hours: 40h/week",
        "✈️ Vacation days: 25+",
        "👶 Parental leave: 480 days",
        "☕ Fika culture (sacred!)"
    ]
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = "\n" + item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # Right side - EU Average with rounded rectangle
    right_box = comparison_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), Inches(1.5), Inches(4), Inches(5.2)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(120, 120, 120)
    right_box.line.width = Pt(2)
    right_box.line.color.rgb = RGBColor(180, 180, 180)

    right_frame = right_box.text_frame
    right_frame.text = "EU Average 🇪🇺"
    right_frame.paragraphs[0].font.size = Pt(28)
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.color.rgb = RGBColor(220, 220, 220)
    right_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    right_items = [
        "👥 Varies widely",
        "⏰ Work hours: 35-45h/week",
        "✈️ Vacation days: 20-25",
        "👶 Parental leave: varies",
        "☕ Less structured breaks"
    ]
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = "\n" + item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 0

    # 5. Add chart: Cost of Living (enhanced styling)
    print("5. Adding cost of living chart...")
    chart_slide = prs.slides.add_slide(blank_slide_layout)
    chart_slide.background.fill.solid()
    chart_slide.background.fill.fore_color.rgb = RGBColor(250, 252, 255)

    # Title with icon
    title_box = chart_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "💰 Average Monthly Costs (SEK)"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    chart_data = CategoryChartData()
    chart_data.categories = ['Rent (1BR)', 'Groceries', 'Transport', 'Entertainment', 'Internet']
    chart_data.add_series('Stockholm', (12000, 3500, 950, 1500, 350))
    chart_data.add_series('Smaller Cities', (8000, 3000, 800, 1200, 300))

    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5.2)
    chart = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    from pptx.enum.chart import XL_LEGEND_POSITION
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(14)

    # Style the chart with Swedish colors
    try:
        # Color the series with Swedish blue and a lighter variant
        plot = chart.plots[0]
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = RGBColor(0, 106, 167)  # Stockholm - Swedish blue
        plot.series[1].format.fill.solid()
        plot.series[1].format.fill.fore_color.rgb = RGBColor(254, 204, 0)  # Smaller cities - Swedish yellow
    except:
        pass  # If styling fails, chart still displays with default colors

    # 6. Add timeline: Year in Sweden (enhanced with visual timeline)
    print("6. Adding timeline slide...")
    timeline_slide = prs.slides.add_slide(blank_slide_layout)
    timeline_slide.background.fill.solid()
    timeline_slide.background.fill.fore_color.rgb = RGBColor(255, 250, 240)

    # Title
    title_box = timeline_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📅 A Year in Sweden"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Draw vertical timeline line
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = timeline_slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1.5), Inches(1.8),
        Inches(1.5), Inches(6.8)
    )
    connector.line.width = Pt(4)
    connector.line.color.rgb = RGBColor(0, 106, 167)

    # Create timeline events with colors for seasons
    events = [
        {"month": "Jan-Feb", "event": "❄️ Polar nights & skiing", "color": RGBColor(100, 150, 200)},
        {"month": "Mar-Apr", "event": "🌷 Spring awakening", "color": RGBColor(150, 200, 150)},
        {"month": "May-Jun", "event": "🌞 Midnight sun begins", "color": RGBColor(255, 220, 100)},
        {"month": "Jul-Aug", "event": "☀️ Summer holidays & festivals", "color": RGBColor(254, 204, 0)},
        {"month": "Sep-Oct", "event": "🍂 Autumn colors & coziness", "color": RGBColor(220, 140, 60)},
        {"month": "Nov-Dec", "event": "🎄 Christmas markets & lights", "color": RGBColor(180, 50, 50)}
    ]

    y_start = 1.6
    for i, evt in enumerate(events):
        # Timeline dot
        dot = timeline_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.3), Inches(y_start + i * 0.85),
            Inches(0.4), Inches(0.4)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = evt["color"]
        dot.line.color.rgb = RGBColor(255, 255, 255)
        dot.line.width = Pt(3)

        # Month box with rounded rectangle
        month_box = timeline_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.2), Inches(y_start + i * 0.85),
            Inches(1.8), Inches(0.5)
        )
        month_box.fill.solid()
        month_box.fill.fore_color.rgb = evt["color"]
        month_box.line.fill.background()
        month_frame = month_box.text_frame
        month_frame.text = evt["month"]
        month_frame.paragraphs[0].font.size = Pt(14)
        month_frame.paragraphs[0].font.bold = True
        month_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        month_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        month_frame.vertical_anchor = 1  # Middle

        # Event text with background
        event_box = timeline_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.3), Inches(y_start + i * 0.85),
            Inches(5), Inches(0.5)
        )
        event_box.fill.solid()
        event_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        event_box.line.color.rgb = evt["color"]
        event_box.line.width = Pt(2)
        event_frame = event_box.text_frame
        event_frame.text = evt["event"]
        event_frame.paragraphs[0].font.size = Pt(15)
        event_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
        event_frame.vertical_anchor = 1  # Middle

    # 7. Add table: Top Cities (styled)
    print("7. Adding table slide...")
    table_slide = prs.slides.add_slide(blank_slide_layout)
    table_slide.background.fill.solid()
    table_slide.background.fill.fore_color.rgb = RGBColor(248, 252, 255)

    # Title
    title_box = table_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🏙️ Top Swedish Cities to Live In"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(8.4)
    height = Inches(5)

    table = table_slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set headers with gradient-like colors
    headers = ['City', 'Population', 'Known For', 'Vibe']
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 106, 167)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER

    # Add data with alternating row colors
    data = [
        ['🏛️ Stockholm', '975K', 'Capital, Islands', 'Cosmopolitan'],
        ['⚓ Gothenburg', '580K', 'Coast, Food Scene', 'Laid-back'],
        ['🎨 Malmö', '345K', 'Diversity, Design', 'International'],
        ['📚 Uppsala', '175K', 'University, History', 'Academic'],
        ['🔬 Lund', '125K', 'Science, Innovation', 'Student City']
    ]

    alternating_colors = [
        RGBColor(255, 255, 255),  # White
        RGBColor(240, 248, 255),  # Light blue
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_value
            cell.fill.solid()
            cell.fill.fore_color.rgb = alternating_colors[row_idx % 2]
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = RGBColor(50, 50, 50)
            if col_idx == 0:  # Make city names bold
                paragraph.font.bold = True

    # 8. Add flowchart: Moving to Sweden Process (enhanced with colors)
    print("8. Adding flowchart...")
    flowchart_slide = prs.slides.add_slide(blank_slide_layout)
    flowchart_slide.background.fill.solid()
    flowchart_slide.background.fill.fore_color.rgb = RGBColor(250, 255, 250)

    # Title
    title_box = flowchart_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "🛫 Moving to Sweden: The Process"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    steps = [
        {"text": "Get job offer or admission", "y": 1.5, "color": RGBColor(0, 106, 167)},
        {"text": "Apply for residence permit", "y": 2.5, "color": RGBColor(30, 130, 180)},
        {"text": "Register with Skatteverket", "y": 3.5, "color": RGBColor(60, 150, 195)},
        {"text": "Get Swedish ID (personnummer)", "y": 4.5, "color": RGBColor(90, 170, 210)},
        {"text": "Enjoy Swedish life! 🎉", "y": 5.5, "color": RGBColor(254, 204, 0)}
    ]

    for i, step in enumerate(steps):
        # Step number circle
        num_circle = flowchart_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(2.2), Inches(step["y"] + 0.15),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = RGBColor(254, 204, 0) if i < len(steps) - 1 else RGBColor(0, 106, 167)
        num_circle.line.fill.background()
        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        num_text.paragraphs[0].font.size = Pt(18)
        num_text.paragraphs[0].font.bold = True
        num_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) if i < len(steps) - 1 else RGBColor(254, 204, 0)
        num_text.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Step box with rounded corners
        shape = flowchart_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2.9), Inches(step["y"]),
            Inches(5), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = step["color"]
        shape.line.fill.background()
        shape.text_frame.text = step["text"]
        shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        shape.text_frame.paragraphs[0].font.size = Pt(16)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.vertical_anchor = 1  # Middle

        # Add arrow connector to next step
        if i < len(steps) - 1:
            connector = flowchart_slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(5.4), Inches(step["y"] + 0.8),
                Inches(5.4), Inches(steps[i+1]["y"])
            )
            connector.line.width = Pt(4)
            connector.line.color.rgb = step["color"]
            # Add arrowhead
            connector.line.dash_style = 1  # Solid line

    # 9. Add content slide with Swedish traditions (enhanced with icons)
    print("9. Adding content slide with Swedish traditions...")
    content_slide = prs.slides.add_slide(blank_slide_layout)
    content_slide.background.fill.solid()
    content_slide.background.fill.fore_color.rgb = RGBColor(255, 252, 245)

    # Title
    title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "💛 Swedish Traditions You'll Love"
    title_frame.paragraphs[0].font.size = Pt(38)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    traditions = [
        {"icon": "☕", "title": "Fika", "desc": "Coffee break culture - sacred daily ritual!", "color": RGBColor(139, 69, 19)},
        {"icon": "🌸", "title": "Midsummer", "desc": "Dancing around maypoles in June", "color": RGBColor(255, 192, 203)},
        {"icon": "🦞", "title": "Crayfish Parties", "desc": "August tradition with paper hats & songs", "color": RGBColor(255, 99, 71)},
        {"icon": "👸", "title": "Lucia", "desc": "December 13th - candlelit procession", "color": RGBColor(255, 215, 0)},
        {"icon": "🌲", "title": "Allemansrätten", "desc": "Freedom to roam nature responsibly", "color": RGBColor(34, 139, 34)}
    ]

    y_start = 1.7
    for i, trad in enumerate(traditions):
        # Icon circle
        icon_circle = content_slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.9), Inches(y_start + i * 0.95),
            Inches(0.6), Inches(0.6)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = trad["color"]
        icon_circle.line.fill.background()
        icon_text = icon_circle.text_frame
        icon_text.text = trad["icon"]
        icon_text.paragraphs[0].font.size = Pt(28)
        icon_text.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Content box
        content_box = content_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.8), Inches(y_start + i * 0.95),
            Inches(7.3), Inches(0.7)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        content_box.line.color.rgb = trad["color"]
        content_box.line.width = Pt(2)

        content_text = content_box.text_frame
        content_text.text = f"{trad['title']}"
        content_text.paragraphs[0].font.size = Pt(18)
        content_text.paragraphs[0].font.bold = True
        content_text.paragraphs[0].font.color.rgb = trad["color"]

        # Description
        p = content_text.add_paragraph()
        p.text = trad['desc']
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.level = 0

    # 10. Add QR code for more info (enhanced design)
    print("10. Adding QR code...")
    import qrcode

    qr_slide = prs.slides.add_slide(blank_slide_layout)
    qr_slide.background.fill.solid()
    qr_slide.background.fill.fore_color.rgb = RGBColor(245, 250, 255)

    # Title
    title_box = qr_slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "📱 Want to Learn More?"
    title_frame.paragraphs[0].font.size = Pt(42)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Decorative frame around QR code
    qr_frame = qr_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.2), Inches(1.8),
        Inches(3.6), Inches(3.6)
    )
    qr_frame.fill.solid()
    qr_frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    qr_frame.line.color.rgb = RGBColor(254, 204, 0)
    qr_frame.line.width = Pt(5)

    # Create QR code for the GitHub repo
    qr = qrcode.QRCode(version=1, box_size=10, border=1)
    qr.add_data("https://github.com/EiriniOr/mcp-powerpoint-server")
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")

    temp_path = "/tmp/qr_demo.png"
    img.save(temp_path)

    qr_slide.shapes.add_picture(temp_path, Inches(3.5), Inches(2.1), width=Inches(3), height=Inches(3))
    os.remove(temp_path)

    # Add text below QR with better styling
    text_box = qr_slide.shapes.add_textbox(Inches(1.5), Inches(5.7), Inches(7), Inches(1.2))
    text_frame = text_box.text_frame
    text_frame.text = "Scan to visit the PowerPoint MCP Server"
    text_frame.paragraphs[0].font.size = Pt(22)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = text_frame.add_paragraph()
    p.text = "GitHub Repository"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    # 11. Thank you slide with enhanced design
    print("11. Adding thank you slide with shapes...")

    thank_you_slide = prs.slides.add_slide(blank_slide_layout)
    thank_you_slide.background.fill.solid()
    thank_you_slide.background.fill.fore_color.rgb = RGBColor(0, 106, 167)

    # Add decorative geometric shapes (Nordic design pattern)
    # Corner triangles
    triangle_positions = [(0.5, 0.5), (8.5, 0.5), (0.5, 6.5), (8.5, 6.5)]
    for x, y in triangle_positions:
        triangle = thank_you_slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(x), Inches(y),
            Inches(0.8), Inches(0.8)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = RGBColor(254, 204, 0)
        triangle.fill.transparency = 0.3
        triangle.line.fill.background()

    # Central decorative hexagons
    hex_positions = [(2, 1.5), (4.5, 0.8), (7, 1.5), (2, 5.7), (4.5, 6.4), (7, 5.7)]
    for x, y in hex_positions:
        hexagon = thank_you_slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON,
            Inches(x), Inches(y),
            Inches(0.6), Inches(0.6)
        )
        hexagon.fill.solid()
        hexagon.fill.fore_color.rgb = RGBColor(254, 204, 0)
        hexagon.fill.transparency = 0.5
        hexagon.line.fill.background()

    # Main content box
    content_box = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.2),
        Inches(7), Inches(3.2)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(254, 204, 0)
    content_box.line.fill.background()

    # Swedish thank you
    title_box = thank_you_slide.shapes.add_textbox(Inches(1.7), Inches(2.5), Inches(6.6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Tack så mycket!"
    title_frame.paragraphs[0].font.size = Pt(56)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # English thank you
    subtitle_box = thank_you_slide.shapes.add_textbox(Inches(1.7), Inches(3.5), Inches(6.6), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Thank you very much!"
    subtitle_frame.paragraphs[0].font.size = Pt(30)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Decorative line
    line_shape = thank_you_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(4.5),
        Inches(3), Inches(0.05)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0, 106, 167)
    line_shape.line.fill.background()

    # Footer
    footer_box = thank_you_slide.shapes.add_textbox(Inches(2), Inches(6.3), Inches(6), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Created with PowerPoint MCP Server"
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = footer_frame.add_paragraph()
    p.text = "36 Powerful Tools for AI-Driven Presentations"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(254, 204, 0)
    p.alignment = PP_ALIGN.CENTER

    # Save the presentation
    output_path = os.path.expanduser("~/Downloads/life_in_sweden_demo.pptx")
    prs.save(output_path)

    print(f"\n✅ Demo presentation created successfully!")
    print(f"📍 Saved to: {output_path}")
    print(f"\n📊 Slides created: {len(prs.slides)}")
    print("\nFeatures showcased:")
    print("  ✓ Title and agenda slides")
    print("  ✓ Section breaks with custom colors")
    print("  ✓ Comparison layout")
    print("  ✓ Column chart with data")
    print("  ✓ Timeline visualization")
    print("  ✓ Formatted table")
    print("  ✓ Flowchart with connectors")
    print("  ✓ Bullet point content")
    print("  ✓ QR code generation")
    print("  ✓ Shapes (stars) and custom styling")
    print("  ✓ Custom backgrounds and colors")

if __name__ == "__main__":
    asyncio.run(create_sweden_demo())
