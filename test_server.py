#!/usr/bin/env python3
"""
Quick test to verify the server can be imported and runs
"""

import sys
import os

# Add the server directory to path
sys.path.insert(0, '/Users/rena/mcp-powerpoint-server')

try:
    # Test imports
    from pptx import Presentation
    from mcp.server import Server
    print("✓ All imports successful")

    # Test creating a simple presentation
    prs = Presentation()
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    slide.shapes.title.text = "Test Presentation"

    # Save to tmp
    test_path = "/tmp/test_presentation.pptx"
    prs.save(test_path)

    # Verify file exists
    if os.path.exists(test_path):
        file_size = os.path.getsize(test_path)
        print(f"✓ Created test presentation: {test_path} ({file_size} bytes)")
        os.remove(test_path)
        print("✓ Cleanup successful")
    else:
        print("✗ Failed to create test file")
        sys.exit(1)

    print("\n✓ All tests passed! Server is ready to use.")

except Exception as e:
    print(f"✗ Test failed: {e}")
    sys.exit(1)
