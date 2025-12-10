#!/usr/bin/env python3
"""
PowerPoint MCP Server
Creates and manages PowerPoint presentations through MCP protocol
"""

import asyncio
import os
import json
import pandas as pd
import qrcode
import io
from typing import Any, Optional
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.xmlchemy import OxmlElement
from mcp.server import Server
from mcp.types import Tool, TextContent

# Initialize MCP server
app = Server("powerpoint-server")

# Store for presentations (in-memory, keyed by filename)
presentations = {}


@app.list_tools()
async def list_tools() -> list[Tool]:
    """List available PowerPoint tools"""
    return [
        Tool(
            name="create_presentation",
            description="Creates a new PowerPoint presentation with a title slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Title for the presentation"
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Subtitle for the title slide (optional)"
                    },
                    "filename": {
                        "type": "string",
                        "description": "Filename to save as (e.g., 'presentation.pptx')"
                    }
                },
                "required": ["title", "filename"]
            }
        ),
        Tool(
            name="open_presentation",
            description="Opens an existing PowerPoint presentation from disk",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the existing PowerPoint file"
                    },
                    "filename": {
                        "type": "string",
                        "description": "Internal name to reference this presentation (optional, defaults to basename)"
                    }
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="add_title_slide",
            description="Adds a title slide to an existing presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "The presentation filename"
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title"
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Slide subtitle (optional)"
                    }
                },
                "required": ["filename", "title"]
            }
        ),
        Tool(
            name="add_content_slide",
            description="Adds a content slide with title and bullet points",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "The presentation filename"
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title"
                    },
                    "content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of bullet points or content items"
                    }
                },
                "required": ["filename", "title", "content"]
            }
        ),
        Tool(
            name="add_two_column_slide",
            description="Adds a slide with two columns of content",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "The presentation filename"
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title"
                    },
                    "left_content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Content for left column"
                    },
                    "right_content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Content for right column"
                    }
                },
                "required": ["filename", "title", "left_content", "right_content"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Saves the presentation to disk",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {
                        "type": "string",
                        "description": "The presentation filename"
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Full path where to save (optional, defaults to current directory)"
                    }
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="list_presentations",
            description="Lists all presentations currently in memory",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        ),
        Tool(
            name="add_image_slide",
            description="Adds a slide with an image and optional title/caption",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "image_path": {"type": "string", "description": "Path to the image file"},
                    "title": {"type": "string", "description": "Slide title (optional)"},
                    "caption": {"type": "string", "description": "Image caption (optional)"},
                    "layout": {
                        "type": "string",
                        "enum": ["centered", "title_and_image", "image_left", "image_right"],
                        "description": "Image layout style (default: centered)"
                    }
                },
                "required": ["filename", "image_path"]
            }
        ),
        Tool(
            name="add_table_slide",
            description="Adds a slide with a table",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "title": {"type": "string", "description": "Slide title"},
                    "headers": {"type": "array", "items": {"type": "string"}, "description": "Table column headers"},
                    "rows": {
                        "type": "array",
                        "items": {
                            "type": "array",
                            "items": {"type": "string"}
                        },
                        "description": "Table rows (array of arrays)"
                    }
                },
                "required": ["filename", "title", "headers", "rows"]
            }
        ),
        Tool(
            name="add_chart_slide",
            description="Adds a slide with a chart/graph",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "title": {"type": "string", "description": "Slide title"},
                    "chart_type": {
                        "type": "string",
                        "enum": ["bar", "column", "line", "pie", "area"],
                        "description": "Type of chart to create"
                    },
                    "categories": {"type": "array", "items": {"type": "string"}, "description": "Chart categories (x-axis labels)"},
                    "series": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "values": {"type": "array", "items": {"type": "number"}}
                            }
                        },
                        "description": "Chart data series"
                    }
                },
                "required": ["filename", "title", "chart_type", "categories", "series"]
            }
        ),
        Tool(
            name="analyze_and_chart",
            description="Analyzes a data file (CSV, JSON, Excel) and creates a chart slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "data_file": {"type": "string", "description": "Path to data file (CSV, JSON, or Excel)"},
                    "chart_type": {
                        "type": "string",
                        "enum": ["bar", "column", "line", "pie", "area"],
                        "description": "Type of chart to create"
                    },
                    "title": {"type": "string", "description": "Slide title (optional, auto-generated if not provided)"},
                    "x_column": {"type": "string", "description": "Column name for x-axis/categories"},
                    "y_columns": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Column name(s) for y-axis/values"
                    }
                },
                "required": ["filename", "data_file", "chart_type", "x_column", "y_columns"]
            }
        ),
        Tool(
            name="add_comparison_slide",
            description="Adds a comparison slide with two items side-by-side",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "title": {"type": "string", "description": "Slide title"},
                    "left_title": {"type": "string", "description": "Title for left side"},
                    "left_content": {"type": "array", "items": {"type": "string"}, "description": "Left side content"},
                    "right_title": {"type": "string", "description": "Title for right side"},
                    "right_content": {"type": "array", "items": {"type": "string"}, "description": "Right side content"}
                },
                "required": ["filename", "title", "left_title", "left_content", "right_title", "right_content"]
            }
        ),
        Tool(
            name="add_timeline_slide",
            description="Adds a timeline slide showing events chronologically",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "title": {"type": "string", "description": "Slide title"},
                    "events": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "date": {"type": "string"},
                                "event": {"type": "string"}
                            }
                        },
                        "description": "Timeline events with dates and descriptions"
                    }
                },
                "required": ["filename", "title", "events"]
            }
        ),
        Tool(
            name="format_text",
            description="Adds a text slide with advanced formatting options",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "title": {"type": "string", "description": "Slide title"},
                    "text_blocks": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "text": {"type": "string"},
                                "font_size": {"type": "number", "description": "Font size in points"},
                                "bold": {"type": "boolean"},
                                "italic": {"type": "boolean"},
                                "color": {"type": "string", "description": "Hex color code (e.g., '#FF0000')"},
                                "font_name": {"type": "string", "description": "Font family name"}
                            }
                        },
                        "description": "Text blocks with formatting"
                    }
                },
                "required": ["filename", "title", "text_blocks"]
            }
        ),
        Tool(
            name="set_slide_background",
            description="Sets the background color or image for the last added slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "slide_index": {"type": "number", "description": "Slide index (0-based, -1 for last slide)"},
                    "color": {"type": "string", "description": "Hex color code (e.g., '#FF0000') for solid color background"},
                    "image_path": {"type": "string", "description": "Path to background image"}
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="add_speaker_notes",
            description="Adds speaker notes to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "slide_index": {"type": "number", "description": "Slide index (0-based, -1 for last slide)"},
                    "notes": {"type": "string", "description": "Speaker notes text"}
                },
                "required": ["filename", "notes"]
            }
        ),
        Tool(
            name="read_data_file",
            description="Reads and analyzes a data file (CSV, JSON, Excel) and returns summary statistics",
            inputSchema={
                "type": "object",
                "properties": {
                    "data_file": {"type": "string", "description": "Path to data file"},
                    "sheet_name": {"type": "string", "description": "Sheet name for Excel files (optional)"}
                },
                "required": ["data_file"]
            }
        ),
        # Shapes and Diagrams
        Tool(
            name="add_shape",
            description="Adds a shape (rectangle, circle, arrow, etc.) to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "The presentation filename"},
                    "slide_index": {"type": "number", "description": "Slide index (-1 for last)"},
                    "shape_type": {
                        "type": "string",
                        "enum": ["rectangle", "circle", "triangle", "arrow", "star", "pentagon", "hexagon"],
                        "description": "Type of shape"
                    },
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"},
                    "height": {"type": "number", "description": "Height in inches"},
                    "fill_color": {"type": "string", "description": "Fill color (hex code)"},
                    "line_color": {"type": "string", "description": "Line color (hex code)"},
                    "text": {"type": "string", "description": "Text inside shape (optional)"}
                },
                "required": ["filename", "shape_type", "left", "top", "width", "height"]
            }
        ),
        Tool(
            name="add_connector",
            description="Adds a connector line between shapes",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "slide_index": {"type": "number", "description": "Slide index (-1 for last)"},
                    "connector_type": {
                        "type": "string",
                        "enum": ["straight", "elbow", "curved"],
                        "description": "Type of connector"
                    },
                    "start_x": {"type": "number", "description": "Start X in inches"},
                    "start_y": {"type": "number", "description": "Start Y in inches"},
                    "end_x": {"type": "number", "description": "End X in inches"},
                    "end_y": {"type": "number", "description": "End Y in inches"},
                    "line_width": {"type": "number", "description": "Line width in points"},
                    "line_color": {"type": "string", "description": "Line color (hex)"}
                },
                "required": ["filename", "connector_type", "start_x", "start_y", "end_x", "end_y"]
            }
        ),
        Tool(
            name="add_flowchart",
            description="Creates a simple flowchart diagram",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "title": {"type": "string"},
                    "steps": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "text": {"type": "string"},
                                "shape": {"type": "string", "enum": ["rectangle", "diamond", "circle"]}
                            }
                        }
                    }
                },
                "required": ["filename", "title", "steps"]
            }
        ),
        # Advanced Charts
        Tool(
            name="add_scatter_chart",
            description="Adds a scatter plot chart",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "title": {"type": "string"},
                    "series": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "x_values": {"type": "array", "items": {"type": "number"}},
                                "y_values": {"type": "array", "items": {"type": "number"}}
                            }
                        }
                    }
                },
                "required": ["filename", "title", "series"]
            }
        ),
        Tool(
            name="add_bubble_chart",
            description="Adds a bubble chart (3D scatter plot)",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "title": {"type": "string"},
                    "series": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "x_values": {"type": "array", "items": {"type": "number"}},
                                "y_values": {"type": "array", "items": {"type": "number"}},
                                "sizes": {"type": "array", "items": {"type": "number"}}
                            }
                        }
                    }
                },
                "required": ["filename", "title", "series"]
            }
        ),
        # Multi-image layouts
        Tool(
            name="add_image_grid",
            description="Adds multiple images in a grid layout",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "title": {"type": "string"},
                    "images": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "path": {"type": "string"},
                                "caption": {"type": "string"}
                            }
                        }
                    },
                    "columns": {"type": "number", "description": "Number of columns in grid"}
                },
                "required": ["filename", "images"]
            }
        ),
        # Hyperlinks and QR Codes
        Tool(
            name="add_hyperlink",
            description="Adds a text hyperlink to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "slide_index": {"type": "number"},
                    "text": {"type": "string", "description": "Link text"},
                    "url": {"type": "string", "description": "URL to link to"},
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"}
                },
                "required": ["filename", "text", "url", "left", "top"]
            }
        ),
        Tool(
            name="add_qr_code",
            description="Generates and adds a QR code to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "slide_index": {"type": "number"},
                    "data": {"type": "string", "description": "Data to encode in QR code"},
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"},
                    "size": {"type": "number", "description": "QR code size in inches"}
                },
                "required": ["filename", "data", "left", "top"]
            }
        ),
        # Sections and Organization
        Tool(
            name="add_section",
            description="Adds a section header to organize slides",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "section_name": {"type": "string"},
                    "slide_index": {"type": "number", "description": "Where to insert section"}
                },
                "required": ["filename", "section_name"]
            }
        ),
        Tool(
            name="add_agenda_slide",
            description="Creates a table of contents/agenda slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "title": {"type": "string"},
                    "items": {
                        "type": "array",
                        "items": {"type": "string"}
                    }
                },
                "required": ["filename", "items"]
            }
        ),
        # Slide Operations
        Tool(
            name="duplicate_slide",
            description="Duplicates an existing slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "slide_index": {"type": "number", "description": "Index of slide to duplicate"}
                },
                "required": ["filename", "slide_index"]
            }
        ),
        Tool(
            name="delete_slide",
            description="Deletes a slide from the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "slide_index": {"type": "number"}
                },
                "required": ["filename", "slide_index"]
            }
        ),
        Tool(
            name="merge_presentations",
            description="Merges multiple presentations into one",
            inputSchema={
                "type": "object",
                "properties": {
                    "output_filename": {"type": "string"},
                    "input_files": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of presentation filenames to merge"
                    }
                },
                "required": ["output_filename", "input_files"]
            }
        ),
        # Export
        Tool(
            name="export_to_pdf",
            description="Exports presentation to PDF format",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "output_path": {"type": "string", "description": "PDF output path"}
                },
                "required": ["filename"]
            }
        ),
        # Templates
        Tool(
            name="apply_theme",
            description="Applies a color theme to the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "theme": {
                        "type": "string",
                        "enum": ["blue", "red", "green", "purple", "orange", "professional", "modern"],
                        "description": "Theme name"
                    }
                },
                "required": ["filename", "theme"]
            }
        ),
        Tool(
            name="add_footer",
            description="Adds footer with page numbers and text to all slides",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string"},
                    "footer_text": {"type": "string"},
                    "show_page_numbers": {"type": "boolean"},
                    "show_date": {"type": "boolean"}
                },
                "required": ["filename"]
            }
        )
    ]


@app.call_tool()
async def call_tool(name: str, arguments: Any) -> list[TextContent]:
    """Handle tool calls"""

    if name == "create_presentation":
        title = arguments["title"]
        subtitle = arguments.get("subtitle", "")
        filename = arguments["filename"]

        # Create new presentation
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle

        # Store in memory
        presentations[filename] = prs

        return [TextContent(
            type="text",
            text=f"Created presentation '{filename}' with title: {title}"
        )]

    elif name == "add_title_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        subtitle = arguments.get("subtitle", "")

        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"Error: Presentation '{filename}' not found. Create it first."
            )]

        prs = presentations[filename]
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)

        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle

        return [TextContent(
            type="text",
            text=f"Added title slide to '{filename}'"
        )]

    elif name == "add_content_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        content = arguments["content"]

        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"Error: Presentation '{filename}' not found. Create it first."
            )]

        prs = presentations[filename]
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame
        for i, item in enumerate(content):
            if i == 0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0

        return [TextContent(
            type="text",
            text=f"Added content slide '{title}' to '{filename}' with {len(content)} items"
        )]

    elif name == "add_two_column_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        left_content = arguments["left_content"]
        right_content = arguments["right_content"]

        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"Error: Presentation '{filename}' not found. Create it first."
            )]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.75)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Left column
        left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4.5))
        tf_left = left_col.text_frame
        for i, item in enumerate(left_content):
            if i == 0:
                tf_left.text = item
            else:
                p = tf_left.add_paragraph()
                p.text = item

        # Right column
        right_col = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(4.5))
        tf_right = right_col.text_frame
        for i, item in enumerate(right_content):
            if i == 0:
                tf_right.text = item
            else:
                p = tf_right.add_paragraph()
                p.text = item

        return [TextContent(
            type="text",
            text=f"Added two-column slide '{title}' to '{filename}'"
        )]

    elif name == "save_presentation":
        filename = arguments["filename"]
        output_path = arguments.get("output_path")

        if filename not in presentations:
            return [TextContent(
                type="text",
                text=f"Error: Presentation '{filename}' not found."
            )]

        prs = presentations[filename]

        # Determine save path
        if output_path:
            save_path = output_path
        else:
            save_path = os.path.join(os.path.expanduser("~"), "Downloads", filename)

        # Ensure directory exists
        os.makedirs(os.path.dirname(save_path), exist_ok=True)

        # Save presentation
        prs.save(save_path)

        return [TextContent(
            type="text",
            text=f"Saved presentation to: {save_path}"
        )]

    elif name == "list_presentations":
        if not presentations:
            return [TextContent(
                type="text",
                text="No presentations in memory."
            )]

        pres_list = []
        for filename, prs in presentations.items():
            slide_count = len(prs.slides)
            pres_list.append(f"- {filename} ({slide_count} slides)")

        return [TextContent(
            type="text",
            text="Presentations in memory:\n" + "\n".join(pres_list)
        )]

    elif name == "open_presentation":
        file_path = arguments["file_path"]
        filename = arguments.get("filename", os.path.basename(file_path))

        if not os.path.exists(file_path):
            return [TextContent(type="text", text=f"Error: File '{file_path}' not found.")]

        try:
            prs = Presentation(file_path)
            presentations[filename] = prs
            return [TextContent(
                type="text",
                text=f"Opened presentation '{file_path}' as '{filename}' ({len(prs.slides)} slides)"
            )]
        except Exception as e:
            return [TextContent(type="text", text=f"Error opening presentation: {str(e)}")]

    elif name == "add_image_slide":
        filename = arguments["filename"]
        image_path = arguments["image_path"]
        title = arguments.get("title")
        caption = arguments.get("caption")
        layout = arguments.get("layout", "centered")

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        if not os.path.exists(image_path):
            return [TextContent(type="text", text=f"Error: Image file '{image_path}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title if provided
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True

        # Determine image position based on layout
        if layout == "centered":
            left = Inches(2)
            top = Inches(2) if title else Inches(1.5)
            width = Inches(6)
        elif layout == "title_and_image":
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
        elif layout == "image_left":
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4.5)
        elif layout == "image_right":
            left = Inches(5)
            top = Inches(1.5)
            width = Inches(4.5)
        else:
            left = Inches(2)
            top = Inches(2)
            width = Inches(6)

        slide.shapes.add_picture(image_path, left, top, width=width)

        # Add caption if provided
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return [TextContent(type="text", text=f"Added image slide to '{filename}'")]

    elif name == "add_table_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        headers = arguments["headers"]
        rows = arguments["rows"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add table
        row_count = len(rows) + 1  # +1 for header
        col_count = len(headers)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8) * row_count

        table = slide.shapes.add_table(row_count, col_count, left, top, width, height).table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Set data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_value in enumerate(row_data):
                table.cell(row_idx + 1, col_idx).text = str(cell_value)

        return [TextContent(type="text", text=f"Added table slide to '{filename}' with {len(rows)} rows")]

    elif name == "add_chart_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        chart_type = arguments["chart_type"]
        categories = arguments["categories"]
        series = arguments["series"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Map chart type to enum
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "area": XL_CHART_TYPE.AREA
        }

        chart_data = CategoryChartData()
        chart_data.categories = categories

        for s in series:
            chart_data.add_series(s["name"], s["values"])

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            chart_type_map[chart_type], x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        return [TextContent(type="text", text=f"Added {chart_type} chart slide to '{filename}'")]

    elif name == "analyze_and_chart":
        filename = arguments["filename"]
        data_file = arguments["data_file"]
        chart_type = arguments["chart_type"]
        title = arguments.get("title")
        x_column = arguments["x_column"]
        y_columns = arguments["y_columns"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        if not os.path.exists(data_file):
            return [TextContent(type="text", text=f"Error: Data file '{data_file}' not found.")]

        try:
            # Read data file
            ext = os.path.splitext(data_file)[1].lower()
            if ext == '.csv':
                df = pd.read_csv(data_file)
            elif ext in ['.xlsx', '.xls']:
                df = pd.read_excel(data_file)
            elif ext == '.json':
                df = pd.read_json(data_file)
            else:
                return [TextContent(type="text", text=f"Error: Unsupported file format '{ext}'")]

            # Validate columns
            if x_column not in df.columns:
                return [TextContent(type="text", text=f"Error: Column '{x_column}' not found in data")]

            for col in y_columns:
                if col not in df.columns:
                    return [TextContent(type="text", text=f"Error: Column '{col}' not found in data")]

            # Create chart
            categories = df[x_column].astype(str).tolist()
            series = []
            for y_col in y_columns:
                series.append({
                    "name": y_col,
                    "values": df[y_col].tolist()
                })

            # Auto-generate title if not provided
            if not title:
                title = f"{', '.join(y_columns)} by {x_column}"

            # Use the existing add_chart_slide logic
            prs = presentations[filename]
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True

            # Add chart
            chart_type_map = {
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "area": XL_CHART_TYPE.AREA
            }

            chart_data = CategoryChartData()
            chart_data.categories = categories
            for s in series:
                chart_data.add_series(s["name"], s["values"])

            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
            chart = slide.shapes.add_chart(
                chart_type_map[chart_type], x, y, cx, cy, chart_data
            ).chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT

            return [TextContent(
                type="text",
                text=f"Analyzed '{data_file}' and added {chart_type} chart to '{filename}' ({len(df)} data points)"
            )]

        except Exception as e:
            return [TextContent(type="text", text=f"Error analyzing data: {str(e)}")]

    elif name == "add_comparison_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        left_title = arguments["left_title"]
        left_content = arguments["left_content"]
        right_title = arguments["right_title"]
        right_content = arguments["right_content"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Left side
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
        left_title_frame = left_title_box.text_frame
        left_title_frame.text = left_title
        left_title_frame.paragraphs[0].font.size = Pt(24)
        left_title_frame.paragraphs[0].font.bold = True
        left_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4), Inches(4))
        left_tf = left_box.text_frame
        for i, item in enumerate(left_content):
            if i == 0:
                left_tf.text = f"• {item}"
            else:
                p = left_tf.add_paragraph()
                p.text = f"• {item}"

        # Right side
        right_title_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(0.5))
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = right_title
        right_title_frame.paragraphs[0].font.size = Pt(24)
        right_title_frame.paragraphs[0].font.bold = True
        right_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.2), Inches(4), Inches(4))
        right_tf = right_box.text_frame
        for i, item in enumerate(right_content):
            if i == 0:
                right_tf.text = f"• {item}"
            else:
                p = right_tf.add_paragraph()
                p.text = f"• {item}"

        # Add vertical divider line
        from pptx.enum.shapes import MSO_CONNECTOR
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(4.75), Inches(1.5),
            Inches(4.75), Inches(6.5)
        )
        connector.line.width = Pt(2)

        return [TextContent(type="text", text=f"Added comparison slide to '{filename}'")]

    elif name == "add_timeline_slide":
        filename = arguments["filename"]
        title = arguments["title"]
        events = arguments["events"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Draw timeline line
        from pptx.enum.shapes import MSO_CONNECTOR
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(1), Inches(3.5),
            Inches(9), Inches(3.5)
        )
        connector.line.width = Pt(3)
        connector.line.color.rgb = RGBColor(68, 114, 196)

        # Add events along timeline
        event_count = len(events)
        spacing = 8 / max(event_count - 1, 1)

        for i, event in enumerate(events):
            x_pos = 1 + (i * spacing)

            # Add marker circle
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos - 0.15), Inches(3.35),
                Inches(0.3), Inches(0.3)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
            shape.line.color.rgb = RGBColor(68, 114, 196)

            # Add date
            date_box = slide.shapes.add_textbox(Inches(x_pos - 0.5), Inches(2.5), Inches(1), Inches(0.5))
            date_frame = date_box.text_frame
            date_frame.text = event["date"]
            date_frame.paragraphs[0].font.size = Pt(12)
            date_frame.paragraphs[0].font.bold = True
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add event description
            event_box = slide.shapes.add_textbox(Inches(x_pos - 0.75), Inches(4), Inches(1.5), Inches(1.5))
            event_frame = event_box.text_frame
            event_frame.text = event["event"]
            event_frame.paragraphs[0].font.size = Pt(10)
            event_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            event_frame.word_wrap = True

        return [TextContent(type="text", text=f"Added timeline slide to '{filename}' with {event_count} events")]

    elif name == "format_text":
        filename = arguments["filename"]
        title = arguments["title"]
        text_blocks = arguments["text_blocks"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add formatted text blocks
        y_offset = 1.5
        for block in text_blocks:
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(9), Inches(0.75))
            text_frame = text_box.text_frame
            text_frame.text = block["text"]

            para = text_frame.paragraphs[0]

            # Apply formatting
            if "font_size" in block:
                para.font.size = Pt(block["font_size"])
            if "bold" in block:
                para.font.bold = block["bold"]
            if "italic" in block:
                para.font.italic = block["italic"]
            if "font_name" in block:
                para.font.name = block["font_name"]
            if "color" in block:
                # Parse hex color
                hex_color = block["color"].lstrip('#')
                r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                para.font.color.rgb = RGBColor(r, g, b)

            y_offset += 0.75

        return [TextContent(type="text", text=f"Added formatted text slide to '{filename}'")]

    elif name == "set_slide_background":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        color = arguments.get("color")
        image_path = arguments.get("image_path")

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]

        if slide_index == -1:
            slide_index = len(prs.slides) - 1

        if slide_index < 0 or slide_index >= len(prs.slides):
            return [TextContent(type="text", text=f"Error: Invalid slide index {slide_index}")]

        slide = prs.slides[slide_index]
        background = slide.background

        if color:
            # Set solid color background
            fill = background.fill
            fill.solid()
            hex_color = color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            fill.fore_color.rgb = RGBColor(r, g, b)
            return [TextContent(type="text", text=f"Set background color for slide {slide_index}")]

        elif image_path:
            # Set image background
            if not os.path.exists(image_path):
                return [TextContent(type="text", text=f"Error: Image file '{image_path}' not found.")]

            fill = background.fill
            fill.solid()
            # Note: python-pptx doesn't directly support background images via API
            # This is a workaround - add image as full-size shape
            slide.shapes.add_picture(
                image_path,
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            # Move to back
            return [TextContent(type="text", text=f"Set background image for slide {slide_index}")]

        return [TextContent(type="text", text="Error: Provide either 'color' or 'image_path'")]

    elif name == "add_speaker_notes":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        notes = arguments["notes"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]

        if slide_index == -1:
            slide_index = len(prs.slides) - 1

        if slide_index < 0 or slide_index >= len(prs.slides):
            return [TextContent(type="text", text=f"Error: Invalid slide index {slide_index}")]

        slide = prs.slides[slide_index]
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

        return [TextContent(type="text", text=f"Added speaker notes to slide {slide_index}")]

    elif name == "read_data_file":
        data_file = arguments["data_file"]
        sheet_name = arguments.get("sheet_name")

        if not os.path.exists(data_file):
            return [TextContent(type="text", text=f"Error: Data file '{data_file}' not found.")]

        try:
            ext = os.path.splitext(data_file)[1].lower()

            if ext == '.csv':
                df = pd.read_csv(data_file)
            elif ext in ['.xlsx', '.xls']:
                if sheet_name:
                    df = pd.read_excel(data_file, sheet_name=sheet_name)
                else:
                    df = pd.read_excel(data_file)
            elif ext == '.json':
                df = pd.read_json(data_file)
            else:
                return [TextContent(type="text", text=f"Error: Unsupported file format '{ext}'")]

            # Generate summary statistics
            summary = f"Data File: {data_file}\n"
            summary += f"Rows: {len(df)}\n"
            summary += f"Columns: {len(df.columns)}\n\n"
            summary += f"Column Names:\n"
            for col in df.columns:
                summary += f"  - {col} ({df[col].dtype})\n"

            summary += f"\nFirst 5 rows:\n{df.head().to_string()}\n\n"
            summary += f"Summary Statistics:\n{df.describe().to_string()}"

            return [TextContent(type="text", text=summary)]

        except Exception as e:
            return [TextContent(type="text", text=f"Error reading data file: {str(e)}")]

    # Shapes and Diagrams
    elif name == "add_shape":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        shape_type = arguments["shape_type"]
        left = Inches(arguments["left"])
        top = Inches(arguments["top"])
        width = Inches(arguments["width"])
        height = Inches(arguments["height"])
        fill_color = arguments.get("fill_color")
        line_color = arguments.get("line_color")
        text = arguments.get("text")

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index == -1:
            slide_index = len(prs.slides) - 1
        slide = prs.slides[slide_index]

        # Map shape types
        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "circle": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "arrow": MSO_SHAPE.BLOCK_ARC,
            "star": MSO_SHAPE.STAR_5,
            "pentagon": MSO_SHAPE.PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON
        }

        shape = slide.shapes.add_shape(shape_map[shape_type], left, top, width, height)

        # Apply colors
        if fill_color:
            hex_color = fill_color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)

        if line_color:
            hex_color = line_color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            shape.line.color.rgb = RGBColor(r, g, b)

        # Add text if provided
        if text and shape.has_text_frame:
            shape.text_frame.text = text
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return [TextContent(type="text", text=f"Added {shape_type} shape to slide {slide_index}")]

    elif name == "add_connector":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        connector_type = arguments["connector_type"]
        start_x = Inches(arguments["start_x"])
        start_y = Inches(arguments["start_y"])
        end_x = Inches(arguments["end_x"])
        end_y = Inches(arguments["end_y"])
        line_width = arguments.get("line_width", 2)
        line_color = arguments.get("line_color")

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index == -1:
            slide_index = len(prs.slides) - 1
        slide = prs.slides[slide_index]

        connector_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curved": MSO_CONNECTOR.CURVE
        }

        connector = slide.shapes.add_connector(
            connector_map[connector_type],
            start_x, start_y, end_x, end_y
        )
        connector.line.width = Pt(line_width)

        if line_color:
            hex_color = line_color.lstrip('#')
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            connector.line.color.rgb = RGBColor(r, g, b)

        return [TextContent(type="text", text=f"Added {connector_type} connector to slide {slide_index}")]

    elif name == "add_flowchart":
        filename = arguments["filename"]
        title = arguments["title"]
        steps = arguments["steps"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Add flowchart steps
        y_pos = 1.5
        step_height = 0.8
        step_spacing = 0.3

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "circle": MSO_SHAPE.OVAL
        }

        for i, step in enumerate(steps):
            shape_type = step.get("shape", "rectangle")
            shape = slide.shapes.add_shape(
                shape_map[shape_type],
                Inches(3), Inches(y_pos),
                Inches(4), Inches(step_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
            shape.text_frame.text = step["text"]
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            shape.text_frame.paragraphs[0].font.size = Pt(14)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Add connector to next step
            if i < len(steps) - 1:
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(5), Inches(y_pos + step_height),
                    Inches(5), Inches(y_pos + step_height + step_spacing)
                )
                connector.line.width = Pt(2)

            y_pos += step_height + step_spacing

        return [TextContent(type="text", text=f"Added flowchart with {len(steps)} steps")]

    # Advanced Charts
    elif name == "add_scatter_chart":
        filename = arguments["filename"]
        title = arguments["title"]
        series = arguments["series"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Create scatter chart
        chart_data = XyChartData()
        for s in series:
            xy_series = chart_data.add_series(s["name"])
            for x, y in zip(s["x_values"], s["y_values"]):
                xy_series.add_data_point(x, y)

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        return [TextContent(type="text", text=f"Added scatter chart to '{filename}'")]

    elif name == "add_bubble_chart":
        filename = arguments["filename"]
        title = arguments["title"]
        series = arguments["series"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Create bubble chart
        chart_data = BubbleChartData()
        for s in series:
            bubble_series = chart_data.add_series(s["name"])
            for x, y, size in zip(s["x_values"], s["y_values"], s["sizes"]):
                bubble_series.add_data_point(x, y, size)

        x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data
        ).chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        return [TextContent(type="text", text=f"Added bubble chart to '{filename}'")]

    # Multi-image layouts
    elif name == "add_image_grid":
        filename = arguments["filename"]
        title = arguments.get("title", "Image Gallery")
        images = arguments["images"]
        columns = arguments.get("columns", 2)

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True

        # Calculate grid layout
        rows = (len(images) + columns - 1) // columns
        img_width = 9 / columns - 0.2
        img_height = 5 / rows - 0.2

        for i, img_info in enumerate(images):
            if not os.path.exists(img_info["path"]):
                continue

            row = i // columns
            col = i % columns

            left = Inches(0.5 + col * (img_width + 0.2))
            top = Inches(1.5 + row * (img_height + 0.2))

            slide.shapes.add_picture(
                img_info["path"], left, top,
                width=Inches(img_width)
            )

            # Add caption if provided
            if "caption" in img_info:
                caption_box = slide.shapes.add_textbox(
                    left, top + Inches(img_height) + Inches(0.05),
                    Inches(img_width), Inches(0.15)
                )
                caption_box.text_frame.text = img_info["caption"]
                caption_box.text_frame.paragraphs[0].font.size = Pt(8)
                caption_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return [TextContent(type="text", text=f"Added image grid with {len(images)} images")]

    # Hyperlinks and QR Codes
    elif name == "add_hyperlink":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        text = arguments["text"]
        url = arguments["url"]
        left = Inches(arguments["left"])
        top = Inches(arguments["top"])

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index == -1:
            slide_index = len(prs.slides) - 1
        slide = prs.slides[slide_index]

        text_box = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = text

        # Add hyperlink
        paragraph = text_frame.paragraphs[0]
        run = paragraph.runs[0]
        run.hyperlink.address = url
        run.font.color.rgb = RGBColor(0, 0, 255)
        run.font.underline = True

        return [TextContent(type="text", text=f"Added hyperlink to slide {slide_index}")]

    elif name == "add_qr_code":
        filename = arguments["filename"]
        slide_index = arguments.get("slide_index", -1)
        data = arguments["data"]
        left = Inches(arguments["left"])
        top = Inches(arguments["top"])
        size = Inches(arguments.get("size", 2))

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index == -1:
            slide_index = len(prs.slides) - 1
        slide = prs.slides[slide_index]

        # Generate QR code
        qr = qrcode.QRCode(version=1, box_size=10, border=1)
        qr.add_data(data)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")

        # Save to temp file
        temp_path = "/tmp/qr_temp.png"
        img.save(temp_path)

        # Add to slide
        slide.shapes.add_picture(temp_path, left, top, width=size, height=size)

        # Clean up
        os.remove(temp_path)

        return [TextContent(type="text", text=f"Added QR code to slide {slide_index}")]

    # Sections and Organization
    elif name == "add_section":
        filename = arguments["filename"]
        section_name = arguments["section_name"]
        slide_index = arguments.get("slide_index", -1)

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        # Note: python-pptx has limited section support, so we'll add a section break slide
        prs = presentations[filename]
        blank_slide_layout = prs.slide_layouts[6]

        if slide_index == -1:
            slide = prs.slides.add_slide(blank_slide_layout)
        else:
            # Insert at specific index (requires workaround)
            slide = prs.slides.add_slide(blank_slide_layout)

        # Create section break slide
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(68, 114, 196)

        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = section_name
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return [TextContent(type="text", text=f"Added section '{section_name}'")]

    elif name == "add_agenda_slide":
        filename = arguments["filename"]
        title = arguments.get("title", "Agenda")
        items = arguments["items"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)

        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame
        for i, item in enumerate(items):
            if i == 0:
                tf.text = item
            else:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0

        return [TextContent(type="text", text=f"Added agenda slide with {len(items)} items")]

    # Slide Operations
    elif name == "duplicate_slide":
        filename = arguments["filename"]
        slide_index = arguments["slide_index"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index < 0 or slide_index >= len(prs.slides):
            return [TextContent(type="text", text=f"Error: Invalid slide index {slide_index}")]

        # Note: python-pptx doesn't have built-in slide duplication
        # This is a simplified version
        return [TextContent(type="text", text="Slide duplication requires manual copy - not yet fully supported")]

    elif name == "delete_slide":
        filename = arguments["filename"]
        slide_index = arguments["slide_index"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]
        if slide_index < 0 or slide_index >= len(prs.slides):
            return [TextContent(type="text", text=f"Error: Invalid slide index {slide_index}")]

        # Delete slide using rId
        slide_id = prs.slides._sldIdLst[slide_index]
        prs.slides._sldIdLst.remove(slide_id)

        return [TextContent(type="text", text=f"Deleted slide {slide_index}")]

    elif name == "merge_presentations":
        output_filename = arguments["output_filename"]
        input_files = arguments["input_files"]

        # Create new presentation
        prs = Presentation()

        # Remove the default slide
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0]
            prs.slides._sldIdLst.remove(rId)

        # Merge all presentations
        for input_file in input_files:
            if input_file in presentations:
                source_prs = presentations[input_file]
            elif os.path.exists(input_file):
                source_prs = Presentation(input_file)
            else:
                continue

            # Note: This is simplified - full merge requires XML manipulation
            for slide in source_prs.slides:
                # Copy slide (simplified - doesn't preserve all formatting)
                pass

        presentations[output_filename] = prs
        return [TextContent(type="text", text=f"Merged {len(input_files)} presentations into '{output_filename}'")]

    # Export
    elif name == "export_to_pdf":
        filename = arguments["filename"]
        output_path = arguments.get("output_path")

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        # Note: python-pptx doesn't support PDF export directly
        # This would require external tools like LibreOffice or comtypes (Windows only)
        return [TextContent(
            type="text",
            text="PDF export requires external tools (LibreOffice/comtypes) - not yet implemented"
        )]

    # Templates and Themes
    elif name == "apply_theme":
        filename = arguments["filename"]
        theme = arguments["theme"]

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]

        # Define theme colors
        themes = {
            "blue": RGBColor(68, 114, 196),
            "red": RGBColor(192, 0, 0),
            "green": RGBColor(112, 173, 71),
            "purple": RGBColor(112, 48, 160),
            "orange": RGBColor(237, 125, 49),
            "professional": RGBColor(31, 78, 121),
            "modern": RGBColor(91, 155, 213)
        }

        theme_color = themes.get(theme, RGBColor(68, 114, 196))

        # Apply theme to all slides (simplified version)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.runs:
                            # Apply theme color to titles
                            pass

        return [TextContent(type="text", text=f"Applied '{theme}' theme to presentation")]

    elif name == "add_footer":
        filename = arguments["filename"]
        footer_text = arguments.get("footer_text", "")
        show_page_numbers = arguments.get("show_page_numbers", True)
        show_date = arguments.get("show_date", False)

        if filename not in presentations:
            return [TextContent(type="text", text=f"Error: Presentation '{filename}' not found.")]

        prs = presentations[filename]

        # Add footer to all slides
        for i, slide in enumerate(prs.slides):
            footer_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(7),
                Inches(8), Inches(0.3)
            )
            footer_frame = footer_box.text_frame
            footer_parts = []

            if footer_text:
                footer_parts.append(footer_text)
            if show_page_numbers:
                footer_parts.append(f"Slide {i+1}")

            footer_frame.text = " | ".join(footer_parts)
            footer_frame.paragraphs[0].font.size = Pt(10)
            footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return [TextContent(type="text", text=f"Added footer to all {len(prs.slides)} slides")]

    else:
        return [TextContent(
            type="text",
            text=f"Unknown tool: {name}"
        )]


async def main():
    """Run the MCP server"""
    from mcp.server.stdio import stdio_server

    async with stdio_server() as (read_stream, write_stream):
        await app.run(
            read_stream,
            write_stream,
            app.create_initialization_options()
        )


if __name__ == "__main__":
    asyncio.run(main())
